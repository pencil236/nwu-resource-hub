import io
import re
import zipfile
from functools import lru_cache
from typing import cast


@lru_cache
def _ocr_model():
    try:
        from paddleocr import PaddleOCR
    except ImportError:
        return None
    return PaddleOCR(
        use_doc_orientation_classify=False,
        use_doc_unwarping=False,
        use_textline_orientation=False,
        lang="ch",
    )


def _ocr_image(data: bytes) -> str:
    try:
        import numpy as np
        from PIL import Image
    except ImportError:
        return ""
    ocr = _ocr_model()
    if ocr is None:
        return ""
    image = np.array(Image.open(io.BytesIO(data)).convert("RGB"))
    texts: list[str] = []
    for item in ocr.predict(image):
        payload = getattr(item, "json", {})
        if callable(payload):
            payload = payload()
        result = payload.get("res", payload) if isinstance(payload, dict) else {}
        texts.extend(str(text) for text in result.get("rec_texts", []) if text)
    return "\n".join(texts)


def extract_text(filename: str, data: bytes) -> str:
    suffix = filename.lower().rsplit(".", 1)[-1]
    if suffix == "pdf":
        import fitz

        document = fitz.open(stream=data, filetype="pdf")
        text = "\n".join(cast(str, page.get_text("text")) for page in document)
        if len(text.strip()) < 20:
            ocr_pages = []
            for page in document:
                image = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5)).tobytes("png")
                ocr_pages.append(_ocr_image(image))
            text = "\n".join(ocr_pages)
        return text
    if suffix == "docx":
        from docx import Document

        document = Document(io.BytesIO(data))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    if suffix == "pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        return "\n".join(
            str(getattr(shape, "text", ""))
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    if suffix == "xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in workbook.worksheets:
            lines.append(sheet.title)
            for row in sheet.iter_rows(values_only=True):
                lines.append(" | ".join(str(cell) for cell in row if cell is not None))
        return "\n".join(lines)
    if suffix in {"png", "jpg", "jpeg"}:
        return _ocr_image(data)
    if suffix in {"txt", "md"}:
        return data.decode("utf-8", errors="ignore")
    if zipfile.is_zipfile(io.BytesIO(data)):
        raise ValueError("unsupported Office archive")
    raise ValueError("unsupported file type")


def chunk_text(text: str, size: int = 2400, overlap: int = 300) -> list[str]:
    normalized = re.sub(r"[ \t]+", " ", text).strip()
    if not normalized:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(normalized):
        end = min(start + size, len(normalized))
        if end < len(normalized):
            boundary = normalized.rfind("\n", start, end)
            if boundary > start + size // 2:
                end = boundary
        chunks.append(normalized[start:end].strip())
        if end == len(normalized):
            break
        start = max(end - overlap, start + 1)
    return chunks
